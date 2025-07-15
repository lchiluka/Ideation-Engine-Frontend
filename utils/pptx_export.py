from io import BytesIO
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def build_pptx_from_df(df: pd.DataFrame, out_stream: BytesIO | str, workflow: str = "default") -> None:
    prs = Presentation()

    def setup_cell(
        cell,
        text,
        *,
        font_name: str = "Calibri",
        font_size: float = 10,
        bold: bool = False,
        align=PP_ALIGN.LEFT
    ):
        cell.margin_top = cell.margin_bottom = cell.margin_left = cell.margin_right = Pt(0)

        tf = cell.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True

        tf.clear()
        tf.text = str(text).strip()

        for para in tf.paragraphs:
            para.alignment = align
            for run in para.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = bold

    for _, row in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Title
        title_shape = slide.shapes.title
        title_shape.text = str(row.get("title", "Untitled")).strip()
        for p in title_shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Calibri"
                r.font.size = Pt(28)
                r.font.bold = True

        # Fields (conditional by workflow)
        if workflow == "Cross-Industry Ideation":
            fields = [
                ("Agent",                   row.get("agent", "")),
                ("Description",             row.get("description", "")),
                ("Industry",                row.get("novelty_reasoning", "")),  # renamed
                ("Original Solution",       row.get("original_solution", "")),  # new
                ("Adaptation Challenges",   row.get("adaptation_challenges", "")),  # new
                ("Feasibility",             row.get("feasibility_reasoning", "")),
                ("Validated TRL",           row.get("validated_trl", "")),
                ("Validated TRL reasoning", row.get("validated_trl_reasoning", "")),
                ("Components",              row.get("components", "")),
                ("References",              row.get("references", "")),
            ]
        else:
            fields = [
                ("Agent",                   row.get("agent", "")),
                ("Description",             row.get("description", "")),
                ("Novelty",                 row.get("novelty_reasoning", "")),
                ("Feasibility",             row.get("feasibility_reasoning", "")),
                ("Validated TRL",           row.get("validated_trl", "")),
                ("Validated TRL reasoning", row.get("validated_trl_reasoning", "")),
                ("Components",              row.get("components", "")),
                ("References",              row.get("references", "")),
            ]

        # Table
        n_rows, n_cols = len(fields) + 1, 2
        left, top, width, height = Inches(0.5), Inches(1.5), Inches(9.4), Inches(5.6)
        tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

        tbl.columns[0].width = Inches(1.0)
        tbl.columns[1].width = Inches(8.4)
        for r in tbl.rows:
            r.height = Pt(18)

        # Header
        setup_cell(tbl.cell(0, 0), "Field", font_size=10, bold=True,  align=PP_ALIGN.CENTER)
        setup_cell(tbl.cell(0, 1), "Value", font_size=10, bold=True,  align=PP_ALIGN.CENTER)

        # Rows
        for i, (lbl, val) in enumerate(fields, start=1):
            setup_cell(tbl.cell(i, 0), lbl, font_size=10, bold=False, align=PP_ALIGN.LEFT)
            setup_cell(tbl.cell(i, 1), val, font_size=10, bold=False, align=PP_ALIGN.LEFT)

    # Save presentation
    if hasattr(out_stream, "write"):
        prs.save(out_stream)
    else:
        prs.save(out_stream)
