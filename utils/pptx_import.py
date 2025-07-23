# ---------------------------------------------------------------------------

# utils.pptx_import.py - load concept cards from PPTX slides
# ---------------------------------------------------------------------------
"""Load one or more concept cards from a PPTX deck.

Each slide is expected to contain a title and a two‑column table
with rows ``Field`` | ``Value``.  Field names are lower‑cased and
converted to snake_case so that a slide exported with
``pptx_export.build_pptx_from_df`` round‑trips cleanly.
"""
from io import BytesIO
from typing import IO, Any, Dict, List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def read_concept_cards(stream: IO[bytes]) -> List[Dict[str, Any]]:
    """Return a list of concept dicts (including any slide pictures) extracted from a PPTX file."""
    prs = Presentation(stream)
    cards: List[Dict[str, Any]] = []

    for slide in prs.slides:
        card: Dict[str, Any] = {}

        # title
        if slide.shapes.title:
            card["title"] = slide.shapes.title.text.strip()

        # images: extract each picture shape as a BytesIO with a `.name` and extension
        media_files: List[BytesIO] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                ext = img.ext  # e.g. ‘png’, ‘jpg’
                bio = BytesIO(img.blob)
                bio.name = f"{card.get('title','slide')}.{ext}"
                media_files.append(bio)
        if media_files:
            card["media"] = media_files

        # first table (if any)
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            tbl = shape.table
            rows = list(tbl.rows)[1:]  # skip header
            for row in rows:
                key = row.cells[0].text.strip().lower().replace(" ", "_")
                val = row.cells[1].text.strip()
                card[key] = val
            break

        cards.append(card)

    return cards